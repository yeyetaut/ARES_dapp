from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Background Color (Dark Industrial Grey)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(18, 18, 18)

    title = slide.shapes.title
    title.text = "ARES PROTOCOL"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 128) # Neon Green Accent
    
    subtitle = slide.placeholders[1]
    subtitle.text = "The Autonomous Settlement Layer for Physical Assets\nHKUST BLOCKCHAIN LAB // 2026"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The $600B Authenticity Gap"
    
    # Add "Industrial" visual elements
    for i in range(3):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + i*3.1), Inches(2), Inches(2.8), Inches(4))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 30, 30)
        box.line.color.rgb = RGBColor(255, 69, 0) # Red-Orange Alert
        
        text_frame = box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        if i == 0:
            p.text = "TRUST DEFICIT\n\nSecondary markets rely on centralized 'experts'—creating bottlenecks and high failure rates."
        elif i == 1:
            p.text = "SETTLEMENT LAG\n\nManual escrow takes 7-14 days. Capital is trapped in inefficient legacy rails."
        else:
            p.text = "AGENT INCAPACITY\n\nAI Agents cannot 'touch' physical goods, preventing autonomous arbitrage."

def add_architecture_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technical Architecture: Tri-Layer Trust"
    
    # Create a simple "Model" using shapes
    # 1. Agent Box
    agent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.5), Inches(1.5))
    agent.text = "AUTONOMOUS LAYER\n(AI Agents / ERC-6551)"
    agent.fill.solid()
    agent.fill.fore_color.rgb = RGBColor(0, 120, 215)
    
    # 2. Marketplace Box
    market = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.75), Inches(2.5), Inches(2.5), Inches(1.5))
    market.text = "SETTLEMENT LAYER\n(Escrow / Marketplace)"
    market.fill.solid()
    market.fill.fore_color.rgb = RGBColor(0, 255, 128)
    market.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # 3. DePIN Box
    depin = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.5), Inches(2.5), Inches(1.5))
    depin.text = "PHYSICAL LAYER\n(DePIN / NFC Hardware)"
    depin.fill.solid()
    depin.fill.fore_color.rgb = RGBColor(128, 0, 255)
    
    # Add connectors (arrows)
    slide.shapes.add_connector(1, Inches(3), Inches(3.25), Inches(3.75), Inches(3.25))
    slide.shapes.add_connector(1, Inches(6.25), Inches(3.25), Inches(7), Inches(3.25))

def add_market_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Convergence"
    
    # Add a "Data Model" visualization
    chart_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(4))
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(20, 20, 20)
    chart_box.line.color.rgb = RGBColor(100, 100, 100)
    
    tf = chart_box.text_frame
    p = tf.add_paragraph()
    p.text = "Total Addressable Market (TAM): $15 Trillion (M2M Commerce by 2030)"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Collectibles: $602B (6.4% CAGR)"
    p = tf.add_paragraph()
    p.text = "• DePIN Infra: $3.5T Projected"
    p = tf.add_paragraph()
    p.text = "• Agentic Trades: 15-25% of all e-commerce by 2030"

def add_reputation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Reputation & Staking: The Trust Score"
    
    # Visualizing the Soulbound NFT / Reputation
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2.5), Inches(3), Inches(3))
    circle.fill.gradient()
    circle.text = "AREP\nSoulbound Score"
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    
    desc_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2.5), Inches(4), Inches(3))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(40, 40, 40)
    
    tf = desc_box.text_frame
    p = tf.add_paragraph()
    p.text = "TRUST DYNAMICS:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "+10: Successful Trade"
    p = tf.add_paragraph()
    p.text = "+5: Valid Verification"
    p = tf.add_paragraph()
    p.text = "-20: Malicious Slashing"

def create_pro_presentation():
    prs = Presentation()
    
    # Set slide size to 16:9 for modern look
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs)
    add_problem_slide(prs)
    add_architecture_diagram_slide(prs)
    add_market_slide(prs)
    add_reputation_slide(prs)
    
    # Add a simple closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "JOIN THE AGENTIC REVOLUTION"
    slide.placeholders[1].text = "ARES.PROTOCOL // FUTURE-PROOF COMMERCE"
    
    prs.save('ARES_Pro_PitchDeck.pptx')

if __name__ == "__main__":
    create_pro_presentation()
