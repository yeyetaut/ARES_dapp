from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "ARES Protocol"
    subtitle.text = "Agentic Resell Ecosystem & Settlement\nHKUST Blockchain Lab // May 15, 2026"

    # --- Slide 2: Executive Summary ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Executive Summary"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "ARES is a machine-to-machine protocol for autonomous trade."
    p = tf.add_paragraph()
    p.text = "• AI-Driven: Autonomous agents trading via ERC-6551 TBAs."
    p = tf.add_paragraph()
    p.text = "• Trustless: Physical authenticity verified by DePIN nodes."
    p = tf.add_paragraph()
    p.text = "• Scalable: Bridging high-value physical goods to DeFi."

    # --- Slide 3: Problem Statement ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "The Problem"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Current physical marketplaces are broken:"
    p = tf.add_paragraph()
    p.text = "• Counterfeit Risk: No reliable on-chain verification."
    p = tf.add_paragraph()
    p.text = "• Trade Friction: Centralized bottlenecks and high fees (15%)."
    p = tf.add_paragraph()
    p.text = "• Manual Only: No support for autonomous agentic commerce."

    # --- Slide 4: Proposed Solution ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "ARES Solution: Tri-Layer Trust"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "1. Autonomous Layer: AI Agents + ERC-6551 Smart Wallets."
    p = tf.add_paragraph()
    p.text = "2. Settlement Layer: State-machine Escrow (USDC)."
    p = tf.add_paragraph()
    p.text = "3. Physical Layer: DePIN Hardware + NFC Cryptography."

    # --- Slide 5: Market Opportunity ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Market Opportunity (2026-2030)"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Physical Collectibles: $602B Market by 2026."
    p = tf.add_paragraph()
    p.text = "• AI Agent Economy: $3T - $5T Transaction Value by 2030."
    p = tf.add_paragraph()
    p.text = "• DePIN Infrastructure: $3.5T Projected Market Cap."

    # --- Slide 6: Technical Architecture ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Technical Architecture"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• AgentRegistry: ERC-6551 TBA auto-deployment."
    p = tf.add_paragraph()
    p.text = "• DigitalTwin (ERC-721): NFC-anchored assets."
    p = tf.add_paragraph()
    p.text = "• Verifier: Staking-backed DePIN nodes."
    p = tf.add_paragraph()
    p.text = "• Reputation: Soulbound (AREP) trust scores."

    # --- Slide 7: Validation & Security ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Validation & Security"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• 100% Test Coverage: 80/80 cases passing."
    p = tf.add_paragraph()
    p.text = "• Slashable Stakes: 100+ USDC collateral for verifiers."
    p = tf.add_paragraph()
    p.text = "• Spending Policies: Hardcoded limits on Agent TBAs."

    # --- Slide 8: Financial Projection ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Financial Forecast (3 Year)"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Year 1 (Pilot): $1M MTV // $15K Revenue"
    p = tf.add_paragraph()
    p.text = "• Year 2 (Growth): $10M MTV // $150K Revenue"
    p = tf.add_paragraph()
    p.text = "• Year 3 (Scale): $50M MTV // $750K Revenue"
    p = tf.add_paragraph()
    p.text = "• Business Model: 1.5% marketplace fee + slashing revenue."

    # --- Slide 9: Conclusion ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "The Future of Autonomous Trade"
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "ARES is not just a marketplace; it is the infrastructure for the machine-to-machine economy."
    p = tf.add_paragraph()
    p.text = "Empowering autonomous agents to trade physical value with digital certainty."

    prs.save('ARES_Presentation.pptx')

if __name__ == "__main__":
    create_presentation()
